"""Export a structured plan dashboard to PDF or PowerPoint (TS-191)."""

from __future__ import annotations

import io
from typing import Any


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


def export_plan_dashboard_to_pdf(dashboard: dict[str, Any], opportunity_title: str = "") -> bytes:
    """Render a plan dashboard as a PDF report using ReportLab."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    title = dashboard.get("title", "Tender Plan Dashboard")
    summary = dashboard.get("summary", "")
    sections = dashboard.get("sections", [])

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        title=title,
    )
    styles = getSampleStyleSheet()
    elements: list[Any] = []

    elements.append(Paragraph(title, styles["Title"]))
    if opportunity_title:
        elements.append(Paragraph(f"Opportunity: {opportunity_title}", styles["Heading2"]))
    if summary:
        elements.append(Paragraph(summary, styles["Normal"]))
    elements.append(Spacer(1, 12))

    for section in sections:
        section_title = section.get("title", "Section")
        section_type = section.get("type", "text")
        data = section.get("data", {})

        elements.append(Paragraph(section_title, styles["Heading2"]))

        if section_type == "kpi":
            value = data.get("value", "-")
            unit = data.get("unit", "")
            label = data.get("label", "")
            elements.append(Paragraph(f"{label}: {value} {unit}".strip(), styles["Normal"]))
        elif section_type == "table":
            columns = data.get("columns", [])
            rows = data.get("rows", [])
            if columns and rows:
                table_data = [[c.get("label", c.get("key", "")) for c in columns]]
                for row in rows:
                    table_data.append([str(row.get(c.get("key", ""), "-")) for c in columns])
                table = Table(table_data, repeatRows=1)
                table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563eb")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ]
                    )
                )
                elements.append(table)
        elif section_type == "chart":
            labels = data.get("labels", [])
            datasets = data.get("datasets", [])
            if labels and datasets:
                lines = [f"Chart: {', '.join(labels)}"]
                for ds in datasets:
                    ds_label = ds.get("label", "value")
                    values = ds.get("data", [])
                    lines.append(f"{ds_label}: {', '.join(str(v) for v in values)}")
                elements.append(Paragraph("<br/>".join(lines), styles["Normal"]))
        elif section_type == "mermaid":
            diagram = data.get("diagram", "")
            elements.append(Paragraph(f"Diagram (Mermaid):<br/>{diagram}", styles["Normal"]))
        else:
            content = data.get("content", "")
            elements.append(Paragraph(str(content), styles["Normal"]))

        elements.append(Spacer(1, 12))

    doc.build(elements)
    buf.seek(0)
    return buf.read()


def export_plan_dashboard_to_pptx(
    dashboard: dict[str, Any], opportunity_title: str = ""
) -> bytes:
    """Render a plan dashboard as a PowerPoint deck using python-pptx."""
    if not _pptx_available():
        raise RuntimeError("python-pptx is not installed")

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title = dashboard.get("title", "Tender Plan Dashboard")
    summary = dashboard.get("summary", "")
    sections = dashboard.get("sections", [])

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = opportunity_title or "TenderShield plan dashboard"

    if summary:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Summary"
        slide.placeholders[1].text = summary

    for section in sections:
        section_title = section.get("title", "Section")
        section_type = section.get("type", "text")
        data = section.get("data", {})

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = section_title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        if section_type == "kpi":
            value = data.get("value", "-")
            unit = data.get("unit", "")
            label = data.get("label", "")
            tf.text = f"{label}: {value} {unit}".strip()
        elif section_type == "table":
            columns = data.get("columns", [])
            rows = data.get("rows", [])
            if columns and rows:
                lines = ["\t".join(c.get("label", c.get("key", "")) for c in columns)]
                for row in rows[:20]:
                    lines.append(
                        "\t".join(str(row.get(c.get("key", ""), "-")) for c in columns)
                    )
                tf.text = "\n".join(lines)
            else:
                tf.text = "No table data."
        elif section_type == "chart":
            labels = data.get("labels", [])
            datasets = data.get("datasets", [])
            if labels and datasets:
                lines = [f"Labels: {', '.join(str(label) for label in labels)}"]
                for ds in datasets:
                    ds_label = ds.get("label", "value")
                    values = ds.get("data", [])
                    lines.append(f"{ds_label}: {', '.join(str(v) for v in values)}")
                tf.text = "\n".join(lines)
            else:
                tf.text = "No chart data."
        elif section_type == "mermaid":
            diagram = data.get("diagram", "")
            tf.text = f"Mermaid diagram:\n{diagram}"
        else:
            content = data.get("content", "")
            tf.text = str(content)

        # Keep body readable
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(12)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


class PlanDashboardExporter:
    """Dispatch exporter for `pdf` and `pptx` formats."""

    FORMATS = {"pdf", "pptx"}

    def export(self, dashboard: dict[str, Any], format: str, opportunity_title: str = "") -> bytes:
        if format == "pdf":
            return export_plan_dashboard_to_pdf(dashboard, opportunity_title)
        if format == "pptx":
            return export_plan_dashboard_to_pptx(dashboard, opportunity_title)
        raise ValueError(f"unsupported_format: {format}")
